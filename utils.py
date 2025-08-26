import os
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
from deep_translator import GoogleTranslator
from langdetect import detect, DetectorFactory
from datetime import datetime
import pandas as pd

# Ensure deterministic language detection
DetectorFactory.seed = 0

# --- New Language Detection Function ---
def detect_and_report_language(text):
    """Detects languages in a block of text and reports if multiple are found."""
    if not text.strip():
        return ""
    try:
        # Split text into chunks to check for multiple languages
        chunks = text.split('\n')
        detected_languages = set()
        for chunk in chunks:
            if chunk.strip():
                detected_languages.add(detect(chunk))
        
        if len(detected_languages) > 1:
            return "multi"
        elif len(detected_languages) == 1:
            return detected_languages.pop()
        else:
            return ""
    except Exception as e:
        print(f"Error during language detection: {e}")
        return ""

# --- Helper function for translation ---
def translate_text_chunk(text, dest_lang):
    """Translates a chunk of text."""
    if not text.strip():
        return text
    try:
        return GoogleTranslator(source='auto', target=dest_lang).translate(text)
    except Exception as e:
        print(f"Could not translate chunk: {e}")
        return text # Return original text if translation fails

# --- DOCX Translation ---
def translate_docx(file_path, dest_lang):
    """Translates a .docx file in-place and returns the path to the new file."""
    doc = Document(file_path)
    for para in doc.paragraphs:
        full_text = para.text
        if full_text.strip():
            translated_text = translate_text_chunk(full_text, dest_lang)
            para.clear()
            para.add_run(translated_text)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                cell.text = translate_text_chunk(cell.text, dest_lang)

    new_file_path = os.path.splitext(file_path)[0] + f"_translated_{dest_lang}.docx"
    doc.save(new_file_path)
    return new_file_path

# --- PPTX Translation ---
def translate_pptx(file_path, dest_lang):
    """Translates a .pptx file in-place and returns the path to the new file."""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = translate_text_chunk(run.text, dest_lang)
    
    new_file_path = os.path.splitext(file_path)[0] + f"_translated_{dest_lang}.pptx"
    prs.save(new_file_path)
    return new_file_path

# --- PDF Translation ---
def translate_pdf(file_path, dest_lang):
    """Translates a .pdf file using redaction and returns the path to the new file."""
    doc = fitz.open(file_path)
    for page in doc:
        blocks = page.get_text("dict")["blocks"]
        for b in blocks:
            if b['type'] == 0: # Text block
                for l in b["lines"]:
                    for s in l["spans"]:
                        original_text = s["text"]
                        if original_text.strip():
                            translated_text = translate_text_chunk(original_text, dest_lang)
                            page.add_redact_annot(s["bbox"], fill=(255, 255, 255))
                            page.apply_redactions()
                            page.insert_text(s["bbox"][:2], translated_text, fontsize=s["size"], fontname="helv", color=s["color"])

    new_file_path = os.path.splitext(file_path)[0] + f"_translated_{dest_lang}.pdf"
    doc.save(new_file_path)
    doc.close()
    return new_file_path

# --- Logging ---
def log_activity(activity_type, file_name, source_language, target_language):
    """Logs user activity to a CSV file."""
    log_file = "user_log.csv"
    log_entry = {
        "timestamp": [datetime.now().strftime("%Y-%m-%d %H:%M:%S")],
        "activity_type": [activity_type],
        "file_name": [file_name],
        "source_language": [source_language],
        "target_language": [target_language],
    }
    df = pd.DataFrame(log_entry)
    if os.path.exists(log_file):
        df.to_csv(log_file, mode="a", header=False, index=False)
    else:
        df.to_csv(log_file, index=False)